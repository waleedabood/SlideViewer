"""Export slides to PowerPoint (.pptx) by rendering each as an image."""

import tempfile
from pathlib import Path

from PyQt6.QtCore import QObject, QSize, Qt, QTimer, QUrl, pyqtSignal
from PyQt6.QtWebEngineWidgets import QWebEngineView

from pptx import Presentation
from pptx.util import Emu

from .models import SlideData, StyleSettings
from .renderer import MarkdownRenderer

SLIDE_W_PX = 1920
SLIDE_H_PX = 1080

PPTX_WIDTH = Emu(12192000)   # 13.333" widescreen
PPTX_HEIGHT = Emu(6858000)   # 7.5"


class PptxExporter(QObject):
    """Renders every slide to a PNG via QWebEngineView, then assembles a .pptx."""

    progress = pyqtSignal(int, int)
    finished = pyqtSignal(str)
    error = pyqtSignal(str)

    def __init__(
        self,
        slide_data: SlideData,
        style: StyleSettings,
        renderer: MarkdownRenderer,
        output_path: str,
        base_url: QUrl,
        parent: QObject | None = None,
    ):
        super().__init__(parent)
        self._slide_data = slide_data
        self._style = style
        self._renderer = renderer
        self._output_path = output_path
        self._base_url = base_url
        self._tmp_dir = Path(tempfile.mkdtemp(prefix="slide_export_"))
        self._image_paths: list[Path] = []
        self._current = 0

        self._view = QWebEngineView()
        self._view.setFixedSize(QSize(SLIDE_W_PX, SLIDE_H_PX))
        self._view.setAttribute(
            Qt.WidgetAttribute.WA_DontShowOnScreen, True
        )
        self._view.show()

    def start(self):
        self._current = 0
        self._image_paths.clear()
        self._render_next()

    def _render_next(self):
        total = self._slide_data.total
        if self._current >= total:
            self._assemble_pptx()
            return

        self.progress.emit(self._current + 1, total)

        md = self._slide_data.content_for_render(self._current)
        html = self._renderer.render(
            md,
            self._style,
            presentation=True,
            slide_number=self._current + 1,
            slide_total=total,
        )

        self._view.loadFinished.connect(self._on_load_finished)
        self._view.page().setHtml(html, self._base_url)

    def _on_load_finished(self, ok: bool):
        self._view.loadFinished.disconnect(self._on_load_finished)
        if not ok:
            self.error.emit(f"Failed to render slide {self._current + 1}")
            return
        QTimer.singleShot(250, self._capture)

    def _capture(self):
        pixmap = self._view.grab(self._view.rect())
        path = self._tmp_dir / f"slide_{self._current:04d}.png"
        pixmap.save(str(path), "PNG")
        self._image_paths.append(path)
        self._current += 1
        self._render_next()

    def _assemble_pptx(self):
        try:
            prs = Presentation()
            prs.slide_width = PPTX_WIDTH
            prs.slide_height = PPTX_HEIGHT
            blank_layout = prs.slide_layouts[6]

            for i, img_path in enumerate(self._image_paths):
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(
                    str(img_path), Emu(0), Emu(0), PPTX_WIDTH, PPTX_HEIGHT
                )
                notes_text = self._slide_data.note_for(i)
                if notes_text.strip():
                    slide.notes_slide.notes_text_frame.text = notes_text.strip()

            prs.save(self._output_path)
            self._cleanup()
            self.finished.emit(self._output_path)
        except Exception as exc:
            self._cleanup()
            self.error.emit(str(exc))

    def _cleanup(self):
        self._view.close()
        for p in self._image_paths:
            p.unlink(missing_ok=True)
        try:
            self._tmp_dir.rmdir()
        except OSError:
            pass
